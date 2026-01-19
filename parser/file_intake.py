from typing import List, Dict
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import os
import zipfile
import tempfile
import shutil

# Try importing unstructured for advanced PDF parsing
try:
    from unstructured.partition.pdf import partition_pdf
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

def parse_pdf(path: str) -> str:
    """
    Parses PDF using 'unstructured' for advanced table/image OCR extraction.
    Falls back to 'fitz' (PyMuPDF) if unstructured fails or isn't installed.
    """
    if UNSTRUCTURED_AVAILABLE:
        try:
            # strategy="hi_res" enables layout analysis (tables) and OCR (images)
            # infer_table_structure=True allows extracting the table as HTML
            elements = partition_pdf(
                filename=path,
                strategy="hi_res",
                infer_table_structure=True,
                extract_images_in_pdf=False, # We want the TEXT from images (OCR), not the image files themselves
                chunking_strategy="by_title", # Helps keep semantic sections together
                max_characters=4000,
                new_after_n_chars=3800,
                combine_text_under_n_chars=2000,
            )
            
            all_text = []
            for el in elements:
                # If it's a table, prefer the HTML representation for better LLM understanding
                if el.category == "Table" and el.metadata.text_as_html:
                    all_text.append(f"\n[TABLE]\n{el.metadata.text_as_html}\n[/TABLE]\n")
                else:
                    all_text.append(el.text)
            
            return "\n\n".join(all_text)
            
        except Exception as e:
            print(f"Unstructured parsing failed for {path}: {e}. Falling back to standard PyMuPDF.")
            # Fallthrough to fitz implementation below

    # Fallback / Standard implementation
    doc = fitz.open(path)
    all_text = []
    for page in doc:
        txt = page.get_text()
        if txt:
            all_text.append(txt)
    doc.close()
    return "\n".join(all_text)

def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
    return "\n".join(all_text)

def parse_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

def parse_txt(path: str, encoding="utf-8") -> str:
    with open(path, "r", encoding=encoding, errors="ignore") as f:
        return f.read()

EXT_MAP = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".txt": parse_txt,
}

def parse_file(path: str) -> List[Dict]:
    filename = os.path.basename(path)
    _, ext = os.path.splitext(filename.lower())

    if ext == ".zip":
        return parse_zip(path)

    parser = EXT_MAP.get(ext)
    if not parser:
        return []

    try:
        content = parser(path)
    except Exception as e:
        print(f"Error parsing {filename}: {e}")
        content = ""

    return [{
        "filename": filename,
        "filetype": ext.replace(".", ""),
        "content": content.strip()
    }]

def parse_folder(folder_path: str, recursive: bool=True) -> List[Dict]:
    results = []
    for root, _, files in os.walk(folder_path):
        for fname in files:
            full = os.path.join(root, fname)
            results.extend(parse_file(full))
        if not recursive:
            break
    return results

def parse_zip(path: str) -> List[Dict]:
    temp = tempfile.mkdtemp(prefix="zip_")
    try:
        with zipfile.ZipFile(path, "r") as z:
            z.extractall(temp)
        return parse_folder(temp, recursive=True)
    finally:
        shutil.rmtree(temp, ignore_errors=True)

if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--folder", action="store_true")
    args = ap.parse_args()

    if args.folder:
        items = parse_folder(args.path)
    else:
        items = parse_file(args.path)
    
    for item in items:
        print("-------------------------------------------")
        print(f"Filename: {item['filename']}")
        print(f"Content Preview: {item['content'][:500]}...")